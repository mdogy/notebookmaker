"""Utility functions for processing lectures and generating notebooks."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING, Literal

if TYPE_CHECKING:
    import nbformat

logger = logging.getLogger(__name__)


def process_lecture(
    input_file: Path,
    output_dir: Path,
    analysis_provider: str = "google",
    generation_provider: str | None = None,
    analysis_model: str | None = None,
    generation_model: str | None = None,
    min_priority: int = 5,
    chunk_size: int = 10,
) -> dict[str, Path]:
    """Process a lecture file using two-phase architecture.

    This function implements the complete two-phase workflow:
    - Phase 1: Analyze PDF with vision LLM to extract code sections
    - Phase 2: Generate notebooks from analysis using text/vision LLM

    Args:
        input_file: Path to the PDF file
        output_dir: Directory where notebooks will be saved
        analysis_provider: LLM provider for Phase 1 (default: "google")
        generation_provider: LLM provider for Phase 2 (default: same as analysis)
        analysis_model: Model for Phase 1 (uses provider default if None)
        generation_model: Model for Phase 2 (uses provider default if None)
        min_priority: Minimum section priority to include (default: 5)
        chunk_size: Pages per API call in Phase 1 (default: 10)

    Returns:
        Dictionary with paths to generated files:
            - "analysis": Path to analysis.json
            - "instructor": Path to instructor notebook
            - "student": Path to student notebook

    Raises:
        ValueError: If the input file format is not supported
        FileNotFoundError: If the input file doesn't exist
    """
    from .analysis import analyze_pdf
    from .generation import generate_notebook_from_analysis, load_analysis_from_file

    # Validate file exists
    if not input_file.exists():
        raise FileNotFoundError(f"Input file not found: {input_file}")

    # Only PDFs supported for two-phase workflow
    if input_file.suffix.lower() != ".pdf":
        raise ValueError(
            f"Two-phase workflow only supports PDF files. "
            f"Got: {input_file.suffix}"
        )

    # Use same provider for generation if not specified
    if generation_provider is None:
        generation_provider = analysis_provider

    logger.info(f"Processing lecture: {input_file.name}")
    logger.info(f"  Analysis provider: {analysis_provider}")
    logger.info(f"  Generation provider: {generation_provider}")

    # Ensure output directory exists
    output_dir.mkdir(parents=True, exist_ok=True)

    # Phase 1: Analyze PDF (or load existing analysis)
    analysis_path = input_file.with_suffix(".analysis.json")

    if analysis_path.exists():
        logger.info(f"Found existing analysis: {analysis_path.name}")
        logger.info("  Loading from file (delete to regenerate)")
        analysis = load_analysis_from_file(analysis_path)
    else:
        logger.info("Running Phase 1: Vision analysis")
        analysis = analyze_pdf(
            pdf_path=input_file,
            output_path=analysis_path,
            chunk_size=chunk_size,
            provider=analysis_provider,  # type: ignore[arg-type]
            model=analysis_model,
        )

    # Phase 2: Generate notebooks from analysis
    logger.info("Running Phase 2: Notebook generation")

    base_name = input_file.stem
    instructor_path = output_dir / f"{base_name}_instructor.ipynb"
    student_path = output_dir / f"{base_name}_student.ipynb"

    # Generate instructor notebook
    generate_notebook_from_analysis(
        analysis=analysis,
        output_path=instructor_path,
        notebook_type="instructor",
        min_priority=min_priority,
        provider=generation_provider,
        model=generation_model,
    )

    # Generate student notebook
    generate_notebook_from_analysis(
        analysis=analysis,
        output_path=student_path,
        notebook_type="student",
        min_priority=min_priority,
        provider=generation_provider,
        model=generation_model,
    )

    return {
        "analysis": analysis_path,
        "instructor": instructor_path,
        "student": student_path,
    }


def extract_pdf_content(pdf_path: Path) -> str:
    """Extract text content from a PDF file.

    Args:
        pdf_path: Path to the PDF file

    Returns:
        Extracted text content from all pages, separated by page breaks

    Raises:
        FileNotFoundError: If the PDF file doesn't exist
        ValueError: If the file is not a valid PDF
    """
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")

    if not pdf_path.suffix.lower() == ".pdf":
        raise ValueError(f"File must be a PDF, got: {pdf_path.suffix}")

    try:
        from pypdf import PdfReader

        reader = PdfReader(pdf_path)
        pages_text = []

        for page_num, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text.strip():  # Only include non-empty pages
                pages_text.append(f"--- Page {page_num} ---\n{text}")

        if not pages_text:
            return ""

        return "\n\n".join(pages_text)

    except Exception as e:
        raise ValueError(f"Error reading PDF file {pdf_path}: {e}") from e


def extract_powerpoint_content(pptx_path: Path) -> str:
    """Extract content from a PowerPoint file.

    Args:
        pptx_path: Path to the PowerPoint file

    Returns:
        Extracted text content from all slides, separated by slide breaks

    Raises:
        FileNotFoundError: If the PowerPoint file doesn't exist
        ValueError: If the file is not a valid PowerPoint
    """
    if not pptx_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

    if pptx_path.suffix.lower() not in {".ppt", ".pptx"}:
        raise ValueError(
            f"File must be a PowerPoint file (.ppt or .pptx), got: {pptx_path.suffix}"
        )

    try:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            slide_text.append(f"--- Slide {slide_num} ---")

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if len(slide_text) > 1:  # More than just the header
                slides_text.append("\n".join(slide_text))

        if not slides_text:
            return ""

        return "\n\n".join(slides_text)

    except Exception as e:
        raise ValueError(
            f"Error reading PowerPoint file {pptx_path}: {e}"
        ) from e


def generate_notebook(
    content: str,
    output_path: Path,
    notebook_type: str = "student",
    provider: str = "anthropic",  # Will be validated at runtime
    model: str | None = None,
) -> None:
    """Generate a Jupyter notebook from lecture content using an LLM.

    Args:
        content: Extracted lecture content (from PDF/PowerPoint)
        output_path: Path where the notebook should be saved
        notebook_type: Type of notebook ("student" or "instructor")
        provider: LLM provider to use (default: "anthropic")
        model: Optional model name (uses provider default if None)

    Raises:
        ValueError: If notebook_type is invalid or content is empty
    """
    if notebook_type not in {"student", "instructor"}:
        raise ValueError(
            f"notebook_type must be 'student' or 'instructor', got: {notebook_type}"
        )

    if not content.strip():
        raise ValueError("Content cannot be empty")

    from typing import Literal, cast

    import nbformat

    from notebookmaker.llm import get_llm_config, get_provider

    # Validate provider at runtime
    valid_providers = {"anthropic", "google", "openai", "openrouter"}
    if provider not in valid_providers:
        provider_list = ", ".join(valid_providers)
        raise ValueError(
            f"Invalid provider: {provider}. Must be one of: {provider_list}"
        )

    # Get LLM configuration
    llm_config = get_llm_config()
    if not model:
        model = llm_config.get("model") if llm_config else None
    if not model:
        # Default models by provider
        defaults = {
            "anthropic": "claude-3-5-sonnet-20241022",
            "google": "gemini-2.0-flash-exp",
            "openai": "gpt-4o-mini",
            "openrouter": "anthropic/claude-3-5-sonnet",
        }
        model = defaults.get(provider, "claude-3-5-sonnet-20241022")

    # Create prompt for the LLM
    prompt = _create_notebook_prompt(content, notebook_type)

    # Get LLM provider (safe cast after validation)
    provider_literal = cast(
        Literal["anthropic", "google", "openai", "openrouter"], provider
    )
    llm = get_provider(provider_literal)

    # Generate notebook content from LLM
    from notebookmaker.llm import LLMMessage

    response = llm.generate(
        messages=[LLMMessage(role="user", content=prompt)],
        model=model,
        temperature=0.3,
        max_tokens=4000,
    )

    # Parse LLM response and create notebook
    nb = _parse_llm_response_to_notebook(response.content, notebook_type)

    # Save notebook
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        nbformat.write(nb, f)  # type: ignore[no-untyped-call]


def _load_prompt_fragment(fragment_name: str) -> str:
    """Load a prompt fragment from the prompts/fragments directory.

    Args:
        fragment_name: Name of the fragment file (without .md extension)

    Returns:
        Content of the fragment file
    """
    from pathlib import Path

    # Get the project root (assuming this file is in src/notebookmaker/)
    project_root = Path(__file__).parent.parent.parent
    fragment_path = project_root / "prompts" / "fragments" / f"{fragment_name}.md"

    if not fragment_path.exists():
        raise FileNotFoundError(f"Prompt fragment not found: {fragment_path}")

    with open(fragment_path, encoding="utf-8") as f:
        content = f.read()

    # Remove the markdown header if present (e.g., "# Role")
    lines = content.strip().split("\n")
    if lines and lines[0].startswith("#"):
        # Skip header and blank line after it
        content = "\n".join(lines[2:] if len(lines) > 2 else lines[1:])

    return content.strip()


def _create_notebook_prompt(content: str, notebook_type: str) -> str:
    """Create a prompt for the LLM to generate notebook content.

    Args:
        content: Extracted lecture content
        notebook_type: Type of notebook ("student" or "instructor")

    Returns:
        Formatted prompt for the LLM
    """
    from pathlib import Path

    # Get the project root
    project_root = Path(__file__).parent.parent.parent

    # Load the appropriate template
    if notebook_type == "student":
        template_path = project_root / "prompts" / "student_notebook_template.md"
    else:
        template_path = project_root / "prompts" / "instructor_notebook_template.md"

    if not template_path.exists():
        raise FileNotFoundError(f"Prompt template not found: {template_path}")

    with open(template_path, encoding="utf-8") as f:
        template = f.read()

    # Load all fragments
    fragments = {
        "role": _load_prompt_fragment("role"),
        "notebook_structure": _load_prompt_fragment("notebook_structure"),
        "student_version": _load_prompt_fragment("student_version"),
        "instructor_version": _load_prompt_fragment("instructor_version"),
        "tone": _load_prompt_fragment("tone"),
        "output_format": _load_prompt_fragment("output_format"),
        "content": content,
    }

    # Replace all placeholders in template
    prompt = template
    for key, value in fragments.items():
        prompt = prompt.replace(f"{{{key}}}", value)

    return prompt


def _parse_llm_response_to_notebook(
    llm_response: str, notebook_type: str
) -> nbformat.NotebookNode:
    """Parse LLM percent-formatted Python response into a Jupyter notebook.

    Args:
        llm_response: Percent-formatted Python from LLM (with # %% markers)
        notebook_type: Type of notebook

    Returns:
        NotebookNode object
    """
    import re

    from nbformat.v4 import new_code_cell, new_markdown_cell, new_notebook

    # Extract Python code if wrapped in markdown code block
    content = llm_response.strip()
    code_block_match = re.search(r"```(?:python)?\s*\n(.*?)\n```", content, re.DOTALL)
    if code_block_match:
        content = code_block_match.group(1).strip()

    # Split by cell markers
    lines = content.split("\n")
    cells: list[nbformat.NotebookNode] = []
    current_cell_type: str | None = None
    current_cell_lines: list[str] = []

    def save_current_cell() -> None:
        """Save accumulated lines as a cell."""
        if current_cell_type is None or not current_cell_lines:
            return

        if current_cell_type == "markdown":
            # Remove '# ' prefix from each line
            markdown_lines = []
            for line in current_cell_lines:
                if line.startswith("# "):
                    markdown_lines.append(line[2:])
                elif line.startswith("#"):
                    # Just '#' with nothing after it
                    markdown_lines.append("")
                else:
                    # Line without '# ' prefix (shouldn't happen but handle it)
                    markdown_lines.append(line)
            cell_content = "\n".join(markdown_lines)
            cells.append(new_markdown_cell(cell_content))  # type: ignore[no-untyped-call]
        else:  # code
            cell_content = "\n".join(current_cell_lines)
            cells.append(new_code_cell(cell_content))  # type: ignore[no-untyped-call]

    for line in lines:
        # Check for cell markers
        if line.strip() == "# %% [markdown]":
            save_current_cell()
            current_cell_type = "markdown"
            current_cell_lines = []
        elif line.strip() == "# %%":
            save_current_cell()
            current_cell_type = "code"
            current_cell_lines = []
        else:
            # Add line to current cell
            if current_cell_type is not None:
                current_cell_lines.append(line)

    # Save the last cell
    save_current_cell()

    # If no cells were created, create an error cell
    if not cells:
        error_msg = (
            "# Error\n\n"
            "Failed to parse notebook content.\n\n"
            "Expected percent-formatted Python with `# %%` markers."
        )
        cells.append(new_markdown_cell(error_msg))  # type: ignore[no-untyped-call]

    # Create notebook
    nb = new_notebook(cells=cells)  # type: ignore[no-untyped-call]

    # Add metadata
    nb.metadata.update(
        {
            "kernelspec": {
                "display_name": "Python 3",
                "language": "python",
                "name": "python3",
            },
            "language_info": {
                "name": "python",
                "version": "3.11.0",
            },
        }
    )

    return nb  # type: ignore[no-any-return]
